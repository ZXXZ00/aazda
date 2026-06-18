import os
import sys
import mimetypes
from abc import ABC, abstractmethod

mimetypes.init()
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Dict, List, Optional

import logging
import fitz
import docx
from python_calamine import CalamineWorkbook
from pptx import Presentation

# ---------------------------------------------------------------------------
# Junk extension blacklist — files that should never be indexed
# ---------------------------------------------------------------------------
JUNK_EXACT_EXTS = {
    # Database internals
    ".sqlite", ".sqlite-wal", ".sqlite-shm",
    ".db", ".db-wal", ".db-shm",
    ".sql", ".sql-wal", ".sql-shm",
    # macOS app databases
    ".kgdb", ".kgdb-wal", ".kgdb-shm",
    ".wfindex", ".wfindex-wal", ".wfindex-shm",
    ".mcat", ".mcat-wal", ".mcat-shm",
    ".notifdb", ".notifdb-wal", ".notifdb-shm",
    ".cloudphotodb", ".cloudphotodb-wal", ".cloudphotodb-shm",
    ".musicdb", ".tvdb", ".itdb",
    # App/system data
    ".dat", ".plist",
    ".toc", ".journal",
    ".tagset", ".tagpool",
    # Spotlight / CoreData index files
    ".indexid", ".indexcompactdirectory", ".indexpostings",
    ".indexarrays", ".indexpositions", ".indexids",
    ".indexgroups", ".indexhead", ".indexscores",
    ".indexbigdates", ".indexdirectory", ".indexupdates",
    ".indexstate", ".indextermids", ".indexpositiontable",
    ".shadowindexhead", ".shadowindexgroups",
    ".shadowindexdirectory", ".shadowindextermids",
    ".shadowindexpositiontable", ".shadowindexarrays",
    ".shadowindexcompactdirectory",
    ".directorystorefile",
    ".partitions", ".buckets", ".offsets", ".header",
    # Qdrant / vector DB internal files
    ".ivf-vector-indexes",
    # Build / runtime artifacts
    ".dylib", ".jar", ".jmod", ".jsa", ".jfc",
    ".metallib", ".sym",
    # Lock / state / log internals
    ".lock", ".ckp", ".state", ".tlog",
    # Thumbnails (system-generated)
    ".ithmb", ".thm",
    # Binary blobs
    ".bin", ".blob", ".sst", ".bfc",
    # Misc system
    ".mdplistc", ".bplist",
    ".last", ".mru",
    ".quantizer",
    ".updates",
    # Java security/config
    ".certs", ".security", ".access", ".policy",
    ".p12", ".keystore", ".crt",
}

JUNK_SUFFIX_PATTERNS = ("_toc", "-wal", "-shm")

def is_junk_ext(ext: str) -> bool:
    """Check if an extension should be skipped entirely."""
    if ext in JUNK_EXACT_EXTS:
        return True
    for pattern in JUNK_SUFFIX_PATTERNS:
        if ext.endswith(pattern):
            return True
    if ".ivf-" in ext or ".index" in ext or ".shadowindex" in ext:
        return True
    return False

# ---------------------------------------------------------------------------
# UTF-8 text sniffing for fallback parsing
# ---------------------------------------------------------------------------
_SNIFF_SIZE = 8192

def is_likely_text(filepath: str) -> bool:
    """Read first N bytes and check if content looks like valid UTF-8 text."""
    try:
        with open(filepath, "rb") as f:
            chunk = f.read(_SNIFF_SIZE)
    except (OSError, PermissionError):
        return False
    if not chunk:
        return False
    try:
        chunk.decode("utf-8")
    except UnicodeDecodeError:
        return False
    printable_chars = sum(1 for b in chunk if 32 <= b <= 126 or b in (9, 10, 13))
    return (printable_chars / len(chunk)) > 0.85

# ---------------------------------------------------------------------------

@dataclass
class ParsedDocument:
    path: str
    name: str
    file_type: str
    size: int
    created_at: str
    updated_at: str
    text_content: Optional[str] = None
    ocr_text: Optional[str] = None
    transcript: Optional[str] = None
    dense_vectors: Dict[str, List[float]] = field(default_factory=dict)
    sparse_vectors: Dict[str, Dict[int, float]] = field(default_factory=dict)
    metadata: Dict[str, Any] = field(default_factory=dict)

class BaseParser(ABC):
    @abstractmethod
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        pass

# 1. Plain Text Parser (.txt, .md, source code)
class PlainTextParser(BaseParser):
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        MAX_PREVIEW_SIZE = 100 * 1024  # 100 KB limit for preview mode
        
        if preview_mode and stats.st_size > MAX_PREVIEW_SIZE:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(MAX_PREVIEW_SIZE)
            content += "\n\n--- [Preview truncated due to size limit] ---"
        else:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        
        _, ext = os.path.splitext(path)
        file_type = "text/plain" if ext == ".txt" else f"text/x-{ext[1:]}" if ext else "text/plain"
        
        return ParsedDocument(
            path=os.path.abspath(path),
            name=os.path.basename(path),
            file_type=file_type,
            size=stats.st_size,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=content
        )

# 2. PDF Parser
class PDFParser(BaseParser):
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        doc = fitz.open(path)
        text_pages = []
        num_pages = len(doc)
        
        pages_to_parse = range(num_pages)
        if preview_mode and num_pages > 5:
            pages_to_parse = range(5)
            
        for page_num in pages_to_parse:
            page = doc.load_page(page_num)
            text = page.get_text()
            if text:
                text_pages.append(text)
        
        content = "\n".join(text_pages)
        if preview_mode and num_pages > 5:
            content += "\n\n--- [Preview truncated: showing first 5 pages only] ---"
        
        # Simple metadata extraction from PDF info
        meta_dict = {}
        if doc.metadata:
            for k, v in doc.metadata.items():
                clean_k = k.replace("/", "").strip().lower()
                if isinstance(v, str) and v.strip():
                    meta_dict[clean_k] = v.strip()
                    
        return ParsedDocument(
            path=os.path.abspath(path),
            name=os.path.basename(path),
            file_type="application/pdf",
            size=stats.st_size,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=content,
            metadata=meta_dict
        )

# 3. Word Document Parser (.docx)
class DocxParser(BaseParser):
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        doc = docx.Document(path)
        
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if preview_mode and len(paragraphs) > 100:
            content = "\n".join(paragraphs[:100])
            content += "\n\n--- [Preview truncated: showing first 100 paragraphs only] ---"
        else:
            content = "\n".join(paragraphs)
        
        # Extract metadata if available
        meta_dict = {}
        try:
            core_props = doc.core_properties
            for prop in ["title", "author", "subject", "keywords", "comments"]:
                val = getattr(core_props, prop, None)
                if val:
                    meta_dict[prop] = str(val)
        except Exception:
            pass

        return ParsedDocument(
            path=os.path.abspath(path),
            name=os.path.basename(path),
            file_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            size=stats.st_size,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=content,
            metadata=meta_dict
        )

# 4. Excel Parser (.xlsx)
class XlsxParser(BaseParser):
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        wb = CalamineWorkbook.from_path(path)
        text_rows = []
        for sheet_name in wb.sheet_names:
            row_count = 0
            sheet_data = wb.get_sheet_by_name(sheet_name).to_python()
            for row in sheet_data:
                row_str = " ".join([str(val) for val in row if val is not None and str(val).strip()])
                if row_str.strip():
                    text_rows.append(row_str)
                    row_count += 1
                    if preview_mode and row_count >= 100:
                        text_rows.append("--- [Sheet preview truncated: showing first 100 rows only] ---")
                        break
        
        content = "\n".join(text_rows)
        return ParsedDocument(
            path=os.path.abspath(path),
            name=os.path.basename(path),
            file_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            size=stats.st_size,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=content
        )

# 5. PowerPoint Parser (.pptx)
class PptxParser(BaseParser):
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        prs = Presentation(path)
        text_runs = []
        
        slides_to_parse = prs.slides
        if preview_mode and len(slides_to_parse) > 5:
            slides_to_parse = slides_to_parse[:5]
            
        for slide in slides_to_parse:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text)
        
        content = "\n".join(text_runs)
        if preview_mode and len(prs.slides) > 5:
            content += "\n\n--- [Preview truncated: showing first 5 slides only] ---"
            
        return ParsedDocument(
            path=os.path.abspath(path),
            name=os.path.basename(path),
            file_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            size=stats.st_size,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=content
        )

# Registry Mapping
_PARSERS: Dict[str, BaseParser] = {
    ".txt": PlainTextParser(),
    ".md": PlainTextParser(),
    ".py": PlainTextParser(),
    ".js": PlainTextParser(),
    ".ts": PlainTextParser(),
    ".tsx": PlainTextParser(),
    ".jsx": PlainTextParser(),
    ".json": PlainTextParser(),
    ".html": PlainTextParser(),
    ".css": PlainTextParser(),
    
    ".pdf": PDFParser(),
    ".docx": DocxParser(),
    ".xlsx": XlsxParser(),
    ".pptx": PptxParser(),
}

class DirectoryParser(BaseParser):
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        abs_path = os.path.abspath(path)
        name = os.path.basename(abs_path)
        if not name:
            name = abs_path
            
        return ParsedDocument(
            path=abs_path,
            name=name,
            file_type="inode/directory",
            size=0,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=name
        )

class MetadataOnlyParser(BaseParser):
    """Indexes file name and metadata only — no content extraction."""
    def parse(self, path: str, stats: os.stat_result, preview_mode: bool = False) -> ParsedDocument:
        abs_path = os.path.abspath(path)
        name = os.path.basename(abs_path)
        _, ext = os.path.splitext(name)
        
        # Guess the standard MIME type, falling back to application/extension or octet-stream
        mime_type, _ = mimetypes.guess_type(abs_path)
        if mime_type:
            file_type = mime_type
        else:
            file_type = f"application/{ext[1:].lower()}" if ext else "application/octet-stream"
        
        return ParsedDocument(
            path=abs_path,
            name=name,
            file_type=file_type,
            size=stats.st_size,
            created_at=datetime.fromtimestamp(stats.st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(stats.st_mtime).isoformat(),
            text_content=name  # Index just the filename so it's searchable
        )

def parse_file(path: str, preview_mode: bool = False) -> Optional[ParsedDocument]:
    """Tiered indexing pipeline:
    1. Directories → DirectoryParser
    2. Junk blacklist → skip (returns None)
    3. Dedicated parser → full content parse
    4. UTF-8 text fallback → index as plain text
    5. Metadata-only → index name/path only
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"File {path} does not exist.")
        
    stats = os.stat(path)
    
    # Tier: Directories
    if os.path.isdir(path):
        return DirectoryParser().parse(path, stats, preview_mode)
    
    _, ext = os.path.splitext(path)
    ext = ext.lower()
    
    # Tier: Junk blacklist → skip
    if ext and is_junk_ext(ext):
        return None
        
    # Tier: SVG files -> metadata-only (avoid parsing massive coordinate text)
    if ext == ".svg":
        return MetadataOnlyParser().parse(path, stats, preview_mode)
    
    # Tier: Dedicated parser
    parser = _PARSERS.get(ext)
    if parser is not None:
        try:
            return parser.parse(path, stats, preview_mode)
        except Exception as e:
            print(f"Parser failed for {path} ({e}), trying fallback...", file=sys.stderr)
    
    # Tier: UTF-8 text fallback
    if is_likely_text(path):
        try:
            return PlainTextParser().parse(path, stats, preview_mode)
        except Exception as e:
            print(f"Text fallback failed for {path} ({e})", file=sys.stderr)
    
    # Tier: Metadata-only (binary files)
    return MetadataOnlyParser().parse(path, stats, preview_mode)

if __name__ == "__main__":
    print("Testing parser registry...")
    # Test on some of the test files
    test_pdf = "./test/PptxGenJS-Demo.pptx"
    if os.path.exists(test_pdf):
        doc = parse_file(test_pdf)
        print(f"Parsed {doc.name}: size={doc.size}, type={doc.file_type}")
        print(f"Text content length: {len(doc.text_content or '')}")
        print(f"Metadata keys: {list(doc.metadata.keys())}")
